# file_processor.py
"""
Service for processing uploaded files (images, PDFs, Word, Excel, PowerPoint).
Uses OCR for document extraction and prepares images for Vision models.
"""

import os
import io
import base64
import logging
import mimetypes
from typing import Dict, Any, Optional, Tuple, List
from PIL import Image
import fitz  # PyMuPDF - fast PDF processing

logger = logging.getLogger(__name__)


class FileProcessor:
    """
    Process various file types for chatbot consumption.
    - Images: Resize and encode for Vision APIs
    - PDFs: Extract text using PyMuPDF
    - Office docs: Extract text (Word, Excel, PowerPoint)
    """

    # Supported file types
    SUPPORTED_IMAGES = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
    SUPPORTED_DOCUMENTS = {'.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}

    # Max image dimensions for Vision API
    MAX_IMAGE_SIZE = (2048, 2048)
    MAX_IMAGE_BYTES = 20 * 1024 * 1024  # 20MB

    def __init__(self):
        pass

    @classmethod
    def is_vision_model(cls, model_name: str) -> bool:
        """
        Check if model supports vision/images.

        Uses the database (llm_models) as the source of truth.
        """
        try:
            from db.models.llm_model import LLMModel
            db_model = LLMModel.get_by_model_id(model_name)
            if db_model:
                return db_model.supports_vision
        except Exception as exc:
            logger.warning(f"[FileProcessor] Failed to resolve model '{model_name}' in llm_models: {exc}")
            return False

        logger.warning(f"[FileProcessor] Model '{model_name}' not found in llm_models")
        return False

    @classmethod
    def get_model_info(cls, model_name: str) -> Optional[Dict[str, Any]]:
        """
        Get full model information from database.

        Returns None if model not found in DB.
        """
        try:
            from db.models.llm_model import LLMModel
            db_model = LLMModel.get_by_model_id(model_name)
            if db_model:
                return db_model.to_dict()
        except Exception:
            pass
        return None

    @classmethod
    def get_file_type(cls, filename: str) -> Optional[str]:
        """Determine file type from filename."""
        ext = os.path.splitext(filename)[1].lower()
        if ext in cls.SUPPORTED_IMAGES:
            return 'image'
        elif ext in cls.SUPPORTED_DOCUMENTS:
            return 'document'
        return None

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """Check if file type is supported."""
        return cls.get_file_type(filename) is not None

    def process_file(
        self,
        file_data: bytes,
        filename: str,
        model_name: str = None
    ) -> Dict[str, Any]:
        """
        Process a file and return extracted content.

        Returns:
            {
                'type': 'image' | 'document',
                'filename': str,
                'text_content': str (for documents or non-vision fallback),
                'image_data': str (base64 encoded, for vision models),
                'mime_type': str,
                'page_count': int (for PDFs),
                'error': str (if processing failed)
            }
        """
        file_type = self.get_file_type(filename)

        if not file_type:
            return {
                'type': 'unknown',
                'filename': filename,
                'error': f'Unsupported file type: {filename}'
            }

        try:
            if file_type == 'image':
                return self._process_image(file_data, filename, model_name)
            else:
                return self._process_document(file_data, filename)
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            return {
                'type': file_type,
                'filename': filename,
                'error': str(e)
            }

    def _process_image(
        self,
        file_data: bytes,
        filename: str,
        model_name: str = None
    ) -> Dict[str, Any]:
        """Process image file for Vision API or OCR fallback."""

        ext = os.path.splitext(filename)[1].lower()
        mime_type = mimetypes.guess_type(filename)[0] or 'image/jpeg'

        # Open and optimize image
        image = Image.open(io.BytesIO(file_data))

        # Convert RGBA to RGB if needed (for JPEG)
        if image.mode == 'RGBA' and ext in ['.jpg', '.jpeg']:
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[3])
            image = background

        # Resize if too large
        if image.size[0] > self.MAX_IMAGE_SIZE[0] or image.size[1] > self.MAX_IMAGE_SIZE[1]:
            image.thumbnail(self.MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)
            logger.info(f"Resized image {filename} to {image.size}")

        # Encode to base64
        buffer = io.BytesIO()
        save_format = 'PNG' if ext == '.png' else 'JPEG'
        image.save(buffer, format=save_format, quality=90)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

        result = {
            'type': 'image',
            'filename': filename,
            'mime_type': mime_type,
            'image_data': image_base64,
            'width': image.size[0],
            'height': image.size[1]
        }

        # If model doesn't support vision, try OCR fallback and drop image payload
        if model_name and not self.is_vision_model(model_name):
            ocr_text = self._ocr_image(image)
            if ocr_text:
                result['text_content'] = ocr_text
                result['ocr_used'] = True
            result.pop('image_data', None)

        return result

    def _ocr_image(self, image: Image.Image) -> Optional[str]:
        """
        Perform OCR on image.
        Uses pytesseract if available, otherwise returns None.
        """
        try:
            import pytesseract
            text = pytesseract.image_to_string(image, lang='deu+eng')
            return text.strip() if text.strip() else None
        except ImportError:
            logger.warning("pytesseract not installed, skipping OCR")
            return None
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return None

    def _process_document(
        self,
        file_data: bytes,
        filename: str
    ) -> Dict[str, Any]:
        """Process document files (PDF, Word, Excel, PowerPoint)."""

        ext = os.path.splitext(filename)[1].lower()

        if ext == '.pdf':
            return self._process_pdf(file_data, filename)
        elif ext in ['.doc', '.docx']:
            return self._process_word(file_data, filename)
        elif ext in ['.xls', '.xlsx']:
            return self._process_excel(file_data, filename)
        elif ext in ['.ppt', '.pptx']:
            return self._process_powerpoint(file_data, filename)
        else:
            return {
                'type': 'document',
                'filename': filename,
                'error': f'Unsupported document type: {ext}'
            }

    def _process_pdf(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract text from PDF using PyMuPDF (fast and efficient).
        """
        try:
            doc = fitz.open(stream=file_data, filetype="pdf")

            text_parts = []
            for page_num, page in enumerate(doc, 1):
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"--- Seite {page_num} ---\n{text}")

            page_count = len(doc)
            doc.close()

            full_text = "\n\n".join(text_parts)

            # If no text extracted, might be scanned PDF - needs OCR
            if not full_text.strip():
                full_text = self._ocr_pdf(file_data)
                ocr_used = True
            else:
                ocr_used = False

            return {
                'type': 'document',
                'filename': filename,
                'mime_type': 'application/pdf',
                'text_content': full_text,
                'page_count': page_count,
                'ocr_used': ocr_used
            }

        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            return {
                'type': 'document',
                'filename': filename,
                'error': f'PDF processing failed: {str(e)}'
            }

    def _ocr_pdf(self, file_data: bytes) -> str:
        """OCR a scanned PDF by converting pages to images."""
        try:
            doc = fitz.open(stream=file_data, filetype="pdf")
            text_parts = []

            for page_num, page in enumerate(doc, 1):
                # Render page to image
                mat = fitz.Matrix(2, 2)  # 2x zoom for better OCR
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")

                # OCR the image
                image = Image.open(io.BytesIO(img_data))
                ocr_text = self._ocr_image(image)

                if ocr_text:
                    text_parts.append(f"--- Seite {page_num} (OCR) ---\n{ocr_text}")

            doc.close()
            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"PDF OCR failed: {e}")
            return ""

    def _process_word(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from Word documents."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(file_data))

            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)

            return {
                'type': 'document',
                'filename': filename,
                'mime_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'text_content': '\n\n'.join(text_parts)
            }

        except ImportError:
            return {
                'type': 'document',
                'filename': filename,
                'error': 'python-docx not installed'
            }
        except Exception as e:
            logger.error(f"Word processing failed: {e}")
            return {
                'type': 'document',
                'filename': filename,
                'error': f'Word processing failed: {str(e)}'
            }

    def _process_excel(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from Excel files."""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)

            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"=== Tabelle: {sheet_name} ==="]

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        sheet_text.append(' | '.join(row_values))

                if len(sheet_text) > 1:  # More than just header
                    text_parts.append('\n'.join(sheet_text))

            return {
                'type': 'document',
                'filename': filename,
                'mime_type': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'text_content': '\n\n'.join(text_parts),
                'sheet_count': len(wb.sheetnames)
            }

        except ImportError:
            return {
                'type': 'document',
                'filename': filename,
                'error': 'openpyxl not installed'
            }
        except Exception as e:
            logger.error(f"Excel processing failed: {e}")
            return {
                'type': 'document',
                'filename': filename,
                'error': f'Excel processing failed: {str(e)}'
            }

    def _process_powerpoint(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_data))

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Folie {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_text.append(row_text)

                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))

            return {
                'type': 'document',
                'filename': filename,
                'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'text_content': '\n\n'.join(text_parts),
                'slide_count': len(prs.slides)
            }

        except ImportError:
            return {
                'type': 'document',
                'filename': filename,
                'error': 'python-pptx not installed'
            }
        except Exception as e:
            logger.error(f"PowerPoint processing failed: {e}")
            return {
                'type': 'document',
                'filename': filename,
                'error': f'PowerPoint processing failed: {str(e)}'
            }

    @staticmethod
    def format_for_prompt(processed_files: List[Dict[str, Any]]) -> str:
        """
        Format processed files for inclusion in the prompt.
        """
        if not processed_files:
            return ""

        parts = []
        for file in processed_files:
            if 'error' in file:
                parts.append(f"[Datei: {file['filename']} - Fehler: {file['error']}]")
            elif file.get('text_content'):
                parts.append(f"[Datei: {file['filename']}]\n{file['text_content']}")

        if not parts:
            return ""

        return "\n\n--- Hochgeladene Dokumente ---\n\n" + "\n\n".join(parts)


# Singleton instance
file_processor = FileProcessor()
